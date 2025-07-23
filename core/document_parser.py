import os
from typing import List, Dict, Any
from pathlib import Path

# Document processing imports
import PyPDF2
from pptx import Presentation
from docx import Document
import pandas as pd

class DocumentParser:
    """Handles parsing of multiple document formats"""
    
    def __init__(self):
        self.supported_formats = {'.pdf', '.pptx', '.docx', '.csv', '.txt', '.md'}
    
    def parse_document(self, file_path: str) -> Dict[str, Any]:
        """Parse document based on file extension"""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        if extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {extension}")
        
        try:
            if extension == '.pdf':
                return self._parse_pdf(file_path)
            elif extension == '.pptx':
                return self._parse_pptx(file_path)
            elif extension == '.docx':
                return self._parse_docx(file_path)
            elif extension == '.csv':
                return self._parse_csv(file_path)
            elif extension in ['.txt', '.md']:
                return self._parse_text(file_path)
        except Exception as e:
            return {
                'filename': file_path.name,
                'content': f"Error parsing document: {str(e)}",
                'chunks': [],
                'metadata': {'error': str(e)}
            }
    
    def _parse_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Parse PDF files"""
        content = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                content += f"\n--- Page {page_num + 1} ---\n{text}\n"
        
        chunks = self._create_chunks(content)
        
        return {
            'filename': file_path.name,
            'content': content,
            'chunks': chunks,
            'metadata': {
                'format': 'pdf',
                'pages': len(pdf_reader.pages)
            }
        }
    
    def _parse_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Parse PowerPoint files"""
        content = ""
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides):
            slide_content = f"\n--- Slide {slide_num + 1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_content += shape.text + "\n"
            
            content += slide_content
        
        chunks = self._create_chunks(content)
        
        return {
            'filename': file_path.name,
            'content': content,
            'chunks': chunks,
            'metadata': {
                'format': 'pptx',
                'slides': len(prs.slides)
            }
        }
    
    def _parse_docx(self, file_path: Path) -> Dict[str, Any]:
        """Parse Word documents"""
        doc = Document(file_path)
        content = ""
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content += paragraph.text + "\n"
        
        chunks = self._create_chunks(content)
        
        return {
            'filename': file_path.name,
            'content': content,
            'chunks': chunks,
            'metadata': {
                'format': 'docx',
                'paragraphs': len(doc.paragraphs)
            }
        }
    
    def _parse_csv(self, file_path: Path) -> Dict[str, Any]:
        """Parse CSV files"""
        try:
            df = pd.read_csv(file_path)
            
            # Create readable content from CSV
            content = f"CSV Data Summary:\n"
            content += f"Columns: {', '.join(df.columns)}\n"
            content += f"Rows: {len(df)}\n\n"
            
            # Add first few rows as sample
            content += "Sample Data:\n"
            content += df.head(10).to_string(index=False)
            
            # Create chunks from each row
            chunks = []
            for idx, row in df.iterrows():
                row_text = f"Row {idx + 1}: " + ", ".join([f"{col}: {val}" for col, val in row.items()])
                chunks.append(row_text)
            
            return {
                'filename': file_path.name,
                'content': content,
                'chunks': chunks,
                'metadata': {
                    'format': 'csv',
                    'rows': len(df),
                    'columns': list(df.columns)
                }
            }
        except Exception as e:
            return {
                'filename': file_path.name,
                'content': f"Error reading CSV: {str(e)}",
                'chunks': [],
                'metadata': {'error': str(e)}
            }
    
    def _parse_text(self, file_path: Path) -> Dict[str, Any]:
        """Parse text and markdown files"""
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        chunks = self._create_chunks(content)
        
        return {
            'filename': file_path.name,
            'content': content,
            'chunks': chunks,
            'metadata': {
                'format': file_path.suffix[1:],
                'size': len(content)
            }
        }
    
    def _create_chunks(self, content: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
        """Split content into overlapping chunks"""
        if not content.strip():
            return []
        
        words = content.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size - overlap):
            chunk = ' '.join(words[i:i + chunk_size])
            if chunk.strip():
                chunks.append(chunk.strip())
        
        return chunks 
